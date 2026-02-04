import streamlit as st
from pptx import Presentation
import os
from openai import OpenAI

# =========================
# CONFIG
# =========================
st.set_page_config(page_title="Excel Help Assistant", layout="centered")

# =========================
# OPENAI CLIENT
# =========================
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    st.error("OpenAI API key not found. Please add OPENAI_API_KEY in Streamlit Secrets.")
    st.stop()

client = OpenAI(api_key=api_key)

# =========================
# UI
# =========================
st.title("Excel Help Assistant")
st.write("Ask any question about creating a Table in Excel")

# =========================
# LOAD PPT CONTENT
# =========================
ppt_text = ""
try:
    ppt = Presentation("How to make a Table in Excel.pptx")
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                ppt_text += shape.text + "\n"
except Exception as e:
    st.error("PPT file not found or unreadable.")
    st.stop()

# =========================
# USER INPUT
# =========================
question = st.text_input("Type your question")

# =========================
# AI RESPONSE
# =========================
if question.strip() != "":
    prompt = f"""
You are an Excel expert.

Use the instructions below to answer clearly and step-by-step.

INSTRUCTIONS:
{ppt_text}

USER QUESTION:
{question}
"""

    try:
        response = client.responses.create(
            model="gpt-4o-mini",
            input=prompt
        )

        st.success(response.output_text)

        # Show video help
        st.video("How to make a Table in Excel.mp4")

    except Exception as e:
        st.error("AI request failed. Please check your API key or model access.")
